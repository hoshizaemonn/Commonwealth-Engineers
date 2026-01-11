#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Web運用報告書 テスト生成システム
テスト用に固定日付で実行し、reports/test/に出力します。
"""

import os
import csv
import re
from datetime import datetime
from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- カラー定義 ---
COLOR_BACKGROUND = RGBColor(250, 248, 245)
COLOR_ORANGE = RGBColor(255, 152, 83)
COLOR_DARK_BROWN = RGBColor(101, 67, 33)
COLOR_LIGHT_BROWN = RGBColor(188, 143, 105)
COLOR_GREEN = RGBColor(76, 175, 80)
COLOR_RED = RGBColor(244, 67, 54)

# ==========================================
# テスト用固定日付設定
# ==========================================
REPORT_MONTH = "2025-10"  # 報告対象月
COMP_MONTH = "2025-09"    # 比較対象月

# ==========================================
# CSVデータ読み込み関数
# ==========================================
def load_csv_data(filepath):
    """CSVファイルを読み込む"""
    data = []
    if not os.path.exists(filepath):
        return data
    with open(filepath, 'r', encoding='utf-8-sig') as f:
        reader = csv.DictReader(f)
        for row in reader:
            data.append(row)
    return data

def get_ga4_traffic_data(data_dir):
    """GA4トラフィック獲得データを取得"""
    csv_path = os.path.join(data_dir, "トラフィック獲得_セッションのメインのチャネル_グループ（デフォルト_チャネル_グループ）.csv")
    
    if not os.path.exists(csv_path):
        return {'sessions': {}, 'new_users': 0, 'total_users': 0}
    
    data = load_csv_data(csv_path)
    
    # チャネル別セッション数を集計
    sessions = {}
    new_users = 0
    total_users = 0
    
    for row in data:
        channel = row.get('セッションのメインのチャネル グループ（デフォルト_チャネル_グループ）', '')
        event = row.get('イベント名', '')
        
        if event == 'session_start' and channel:
            try:
                sessions[channel] = int(row.get('セッション', 0))
            except:
                pass
        
        # 新規ユーザー数（first_visitイベントから取得）
        if event == 'first_visit':
            try:
                new_users_value = int(row.get('新規ユーザー数', 0))
                new_users = max(new_users, new_users_value)
            except:
                pass
    
    # イベントデータから総ユーザー数を取得
    events_path = os.path.join(data_dir, "イベント_イベント名.csv")
    if os.path.exists(events_path):
        events_data = load_csv_data(events_path)
        for row in events_data:
            if row.get('イベント名', '') == 'user_engagement':
                try:
                    total_users = int(row.get('総ユーザー数', 0))
                    break
                except:
                    pass
    
    return {
        'sessions': sessions,
        'new_users': new_users,
        'total_users': total_users
    }

def get_ga4_events_data(data_dir):
    """GA4イベントデータを取得"""
    csv_path = os.path.join(data_dir, "イベント_イベント名.csv")
    
    if not os.path.exists(csv_path):
        return {}
    
    data = load_csv_data(csv_path)
    
    events = {}
    total_users = 0
    new_users = 0
    
    for row in data:
        event_name = row.get('イベント名', '')
        if event_name:
            try:
                events[event_name] = int(row.get('イベント数', 0))
                
                if event_name == 'user_engagement':
                    try:
                        total_users = int(row.get('総ユーザー数', 0))
                    except:
                        pass
                
                if event_name == 'first_visit':
                    try:
                        new_users = int(row.get('総ユーザー数', 0))
                    except:
                        pass
            except:
                pass
    
    events['_total_users'] = total_users
    events['_new_users'] = new_users
    
    return events

def get_search_console_data(data_dir):
    """Search Consoleデータを取得"""
    queries_path = os.path.join(data_dir, "クエリ.csv")
    
    result = {
        'total_clicks': 0,
        'total_impressions': 0,
        'avg_ctr': 0.0,
        'avg_position': 0.0,
        'top_queries': []
    }
    
    if os.path.exists(queries_path):
        data = load_csv_data(queries_path)
        total_clicks = 0
        total_impressions = 0
        
        for row in data:
            try:
                clicks = int(row.get('クリック数', 0))
                impressions = int(row.get('表示回数', 0))
                
                total_clicks += clicks
                total_impressions += impressions
            except:
                pass
        
        result['total_clicks'] = total_clicks
        result['total_impressions'] = total_impressions
        result['avg_ctr'] = (total_clicks / total_impressions * 100) if total_impressions > 0 else 0
    
    # デバイスデータから平均順位を計算
    devices_path = os.path.join(data_dir, "デバイス.csv")
    if os.path.exists(devices_path):
        data = load_csv_data(devices_path)
        positions = []
        for row in data:
            try:
                positions.append(float(row.get('掲載順位', 0)))
            except:
                pass
        if positions:
            result['avg_position'] = sum(positions) / len(positions)
    
    return result

def parse_markdown_for_previous_month_data(comparison_year_month):
    """先月レポートから比較対象月のデータを抽出（エラー時は空のデータを返す）"""
    report_path = f"reports/{comparison_year_month}.md"
    
    comparison_data = {}
    previous_previous_data = {}
    
    if os.path.exists(report_path):
        try:
            with open(report_path, 'r', encoding='utf-8') as f:
                content = f.read()
                
                # セッション数を抽出
                pattern = r'\| 総セッション数 \| (\d+) \| (\d+) \|'
                match = re.search(pattern, content)
                if match:
                    previous_previous_data['sessions'] = int(match.group(1))
                    comparison_data['sessions'] = int(match.group(2))
                
                # 問い合わせ件数
                pattern = r'\| 問い合わせ件数 \| (\d+)件 \| (\d+)件 \|'
                match = re.search(pattern, content)
                if match:
                    previous_previous_data['inquiries'] = int(match.group(1))
                    comparison_data['inquiries'] = int(match.group(2))
                
                # CVR
                pattern = r'\| 問い合わせCVR \| ([\d.]+)% \| ([\d.]+)% \|'
                match = re.search(pattern, content)
                if match:
                    previous_previous_data['cvr'] = float(match.group(1))
                    comparison_data['cvr'] = float(match.group(2))
                
                # 自然検索
                pattern = r'\| Organic Search（自然検索） \| (\d+) \| (\d+) \|'
                match = re.search(pattern, content)
                if match:
                    previous_previous_data['organic_search'] = int(match.group(1))
                    comparison_data['organic_search'] = int(match.group(2))
                
                # Direct
                pattern = r'\| Direct（直接流入） \| (\d+) \| (\d+) \|'
                match = re.search(pattern, content)
                if match:
                    previous_previous_data['direct'] = int(match.group(1))
                    comparison_data['direct'] = int(match.group(2))
        except Exception as e:
            print(f"⚠️  レポート読み込みエラー: {e}")
            print(f"   先月のデータなしとして処理を続行します")
    else:
        print(f"⚠️  レポートファイルが見つかりません: {report_path}")
        print(f"   先月のデータなしとして処理を続行します")
    
    return previous_previous_data, comparison_data

# ==========================================
# Markdownファイルの生成機能（強化版）
# ==========================================
def export_to_markdown(report_month, comp_month, report_data, comparison_data, sc_data, comparison_sc_data, output_dir):
    """Markdownファイルを生成（クライアント返信用メール案を含む）"""
    
    # 月の日本語表記
    report_date = datetime.strptime(report_month, "%Y-%m")
    comp_date = datetime.strptime(comp_month, "%Y-%m")
    report_month_str = report_date.strftime("%Y年%m月")
    comp_month_str = comp_date.strftime("%Y年%m月")
    report_month_num = report_date.strftime("%m")
    
    # 良かった点と課題を抽出
    good_points = []
    issues = []
    
    if report_data.get('inquiries_change', 0) > 0:
        good_points.append(f"問い合わせ件数が{report_data.get('inquiries', 0)}件（前月比 +{report_data.get('inquiries_change', 0):.1f}%）に増加")
    
    if report_data.get('cvr_change', 0) > 0:
        good_points.append(f"問い合わせCVRが{report_data.get('cvr', 0):.2f}%（前月比 +{report_data.get('cvr_change', 0):.2f}pt）に改善")
    
    if report_data.get('sessions_change', 0) < 0:
        issues.append(f"セッション数が{report_data.get('sessions', 0)}（前月比 {report_data.get('sessions_change', 0):.1f}%）に減少")
    
    if not good_points:
        good_points.append("サイトの安定稼働を確認")
    if not issues:
        issues.append("引き続きモニタリングが必要")
    
    # 改善提案
    improvement_proposal = """
1. 求人ページ（/entry/）のSEO設定改善
   - タイトル・ディスクリプションの最適化
   - CTR向上を目指す（目標: 3〜5%）

2. コンテンツ拡充
   - プロジェクトページの充実
   - 検索流入の強化
"""
    
    md_content = f"""# {report_month_str} Web運用報告データ

## 1. 主要KPI

| 指標 | {comp_month_str} | {report_month_str} | 前月比 |
| :--- | :--- | :--- | :--- |
| 総セッション数 | {comparison_data.get('sessions', 0)} | {report_data.get('sessions', 0)} | {report_data.get('sessions_change', 0):+.1f}% |
| 問い合わせ件数 | {comparison_data.get('inquiries', 0)}件 | {report_data.get('inquiries', 0)}件 | {report_data.get('inquiries_change', 0):+.1f}% |
| 問い合わせCVR | {comparison_data.get('cvr', 0):.2f}% | {report_data.get('cvr', 0):.2f}% | {report_data.get('cvr_change', 0):+.2f}pt |
| 新規ユーザー数 | - | {report_data.get('new_users', 0)} | - |
| 総ユーザー数 | - | {report_data.get('total_users', 0)} | - |

## 2. 流入経路

| チャネル | {comp_month_str} | {report_month_str} | 前月比 |
| :--- | :--- | :--- | :--- |
| Organic Search | {comparison_data.get('organic_search', 0)} | {report_data.get('organic_search', 0)} | {report_data.get('organic_search_change', 0):+.1f}% |
| Direct | {comparison_data.get('direct', 0)} | {report_data.get('direct', 0)} | {report_data.get('direct_change', 0):+.1f}% |

## 3. Search Console

| 指標 | {comp_month_str} | {report_month_str} | 前月比 |
| :--- | :--- | :--- | :--- |
| 平均CTR | {comparison_sc_data.get('avg_ctr', 0):.2f}% | {sc_data.get('avg_ctr', 0):.2f}% | {(sc_data.get('avg_ctr', 0) - comparison_sc_data.get('avg_ctr', 0)):+.2f}pt |
| 平均掲載順位 | {comparison_sc_data.get('avg_position', 0):.1f}位 | {sc_data.get('avg_position', 0):.1f}位 | {(sc_data.get('avg_position', 0) - comparison_sc_data.get('avg_position', 0)):+.1f}位 |
| 総クリック数 | {comparison_sc_data.get('total_clicks', 0)} | {sc_data.get('total_clicks', 0)} | {sc_data.get('total_clicks', 0) - comparison_sc_data.get('total_clicks', 0):+d} |
| 総表示回数 | {comparison_sc_data.get('total_impressions', 0)} | {sc_data.get('total_impressions', 0)} | {sc_data.get('total_impressions', 0) - comparison_sc_data.get('total_impressions', 0):+d} |

## 4. 特記事項（課題）

- 採用エントリーページ(/entry/)の表示回数: 432回
- 採用エントリーページ(/entry/)のクリック数: 0回
- 課題: 表示されているが選ばれていない（改善パターン1該当）

---

## 【クライアント返信用メール案】

```
件名：【{report_month_num}月分レポート】コモンウェルスエンジニアーズ様 Web運用報告書のご送付

コモンウェルスエンジニアーズ様

お世話になっております。
鈴木です。

{report_month_num}月分のWeb運用報告書を作成いたしましたので、お送りいたします。

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
■ レポート資料
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

【スライド資料】
{{SLIDE_URL}}

【詳細データ（Markdown）】
{{MARKDOWN_URL}}


━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
■ {report_month_num}月のサマリー（前月対比）
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

【良かった点】
✅ {good_points[0] if len(good_points) > 0 else 'サイトの安定稼働を確認'}
✅ {good_points[1] if len(good_points) > 1 else '引き続き改善を継続'}

【課題・懸念点】
⚠️ {issues[0] if len(issues) > 0 else '引き続きモニタリングが必要'}
⚠️ {issues[1] if len(issues) > 1 else '改善施策の継続実施'}


━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
■ 改善提案
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

{improvement_proposal}


━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

ご不明な点やご質問がございましたら、お気軽にお申し付けください。

よろしくお願いいたします。

鈴木
```
"""
    
    # 出力ディレクトリが存在しない場合は作成
    os.makedirs(output_dir, exist_ok=True)
    
    output_path = os.path.join(output_dir, f"{report_month}_レポート.md")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(md_content)
    print(f"✅ Markdownレポートを生成しました: {output_path}")

# ==========================================
# メイン処理
# ==========================================
def main():
    """メイン処理"""
    print("🚀 テスト用PowerPointレポート生成開始")
    print("="*60)
    
    # テスト用固定日付
    report_month = REPORT_MONTH
    comp_month = COMP_MONTH
    
    print(f"📅 報告対象月: {report_month}")
    print(f"📅 比較対象月: {comp_month}")
    
    # 出力ディレクトリを設定
    output_dir = "reports/test"
    os.makedirs(output_dir, exist_ok=True)
    print(f"📁 出力先: {output_dir}")
    
    # データディレクトリ
    data_dir = f"data/{report_month}"
    
    if not os.path.exists(data_dir):
        print(f"⚠️  データディレクトリが見つかりません: {data_dir}")
        print(f"   テスト用のサンプルデータを使用します")
        # サンプルデータを使用
        report_data = {
            'sessions': 1500,
            'inquiries': 24,
            'cvr': 2.88,
            'organic_search': 850,
            'direct': 320,
            'new_users': 1200,
            'total_users': 1400,
            'sessions_change': -5.4,
            'inquiries_change': 33.3,
            'cvr_change': 1.16,
            'organic_search_change': 13.3,
            'direct_change': -36.9
        }
        comparison_data = {
            'sessions': 1047,
            'inquiries': 18,
            'cvr': 1.72,
            'organic_search': 271,
            'direct': 754
        }
        sc_data = {
            'total_clicks': 200,
            'total_impressions': 3500,
            'avg_ctr': 6.02,
            'avg_position': 12.5
        }
        comparison_sc_data = {
            'total_clicks': 168,
            'total_impressions': 3290,
            'avg_ctr': 5.11,
            'avg_position': 8.2
        }
    else:
        # 実際のデータを読み込み
        print(f"📊 データを読み込み中: {data_dir}")
        
        # GA4データ
        ga4_traffic = get_ga4_traffic_data(data_dir)
        ga4_events = get_ga4_events_data(data_dir)
        
        sessions_dict = ga4_traffic.get('sessions', {})
        total_sessions = sum(sessions_dict.values())
        organic_search_sessions = sessions_dict.get('Organic Search', 0)
        direct_sessions = sessions_dict.get('Direct', 0)
        
        new_users = ga4_traffic.get('new_users', 0)
        total_users = ga4_traffic.get('total_users', 0)
        
        if total_users == 0:
            total_users = ga4_events.get('_total_users', 0)
        if new_users == 0:
            new_users = ga4_events.get('_new_users', 0)
        
        form_starts = ga4_events.get('form_start', 0)
        inquiries = ga4_events.get('contact', 0)
        
        cvr = (inquiries / total_sessions * 100) if total_sessions > 0 else 0
        
        # 比較対象月のデータを取得
        previous_previous_data, comparison_data = parse_markdown_for_previous_month_data(comp_month)
        
        comparison_sessions = comparison_data.get('sessions', 0)
        comparison_inquiries = comparison_data.get('inquiries', 0)
        comparison_cvr = comparison_data.get('cvr', 0)
        comparison_organic = comparison_data.get('organic_search', 0)
        comparison_direct = comparison_data.get('direct', 0)
        
        report_data = {
            'sessions': total_sessions,
            'inquiries': inquiries,
            'cvr': cvr,
            'organic_search': organic_search_sessions,
            'direct': direct_sessions,
            'new_users': new_users,
            'total_users': total_users,
            'sessions_change': ((total_sessions - comparison_sessions) / comparison_sessions * 100) if comparison_sessions > 0 else 0,
            'inquiries_change': ((inquiries - comparison_inquiries) / comparison_inquiries * 100) if comparison_inquiries > 0 else 0,
            'cvr_change': cvr - comparison_cvr,
            'organic_search_change': ((organic_search_sessions - comparison_organic) / comparison_organic * 100) if comparison_organic > 0 else 0,
            'direct_change': ((direct_sessions - comparison_direct) / comparison_direct * 100) if comparison_direct > 0 else 0
        }
        
        # Search Consoleデータ
        sc_data = get_search_console_data(data_dir)
        comparison_sc_data = {
            'total_clicks': 0,
            'total_impressions': 0,
            'avg_ctr': 0.0,
            'avg_position': 0.0
        }
    
    # パワーポイントの生成（簡易版）
    print("\n📄 パワーポイントを生成中...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # タイトルスライド（簡易版）
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BACKGROUND
    
    report_date = datetime.strptime(report_month, "%Y-%m")
    month_str = report_date.strftime("%Y年%m月")
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = f"{month_str}度 Web運用報告書"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_DARK_BROWN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf2 = subtitle_box.text_frame
    tf2.text = "コモンウェルスエンジニアーズ様"
    tf2.paragraphs[0].font.size = Pt(24)
    tf2.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    date_str = datetime.now().strftime("%Y年%m月%d日")
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    tf3 = date_box.text_frame
    tf3.text = f"作成日: {date_str}"
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    pptx_name = os.path.join(output_dir, f"{report_month}_レポート.pptx")
    prs.save(pptx_name)
    print(f"✅ パワーポイントを生成しました: {pptx_name}")
    
    # Markdownの生成
    print("\n📝 Markdownレポートを生成中...")
    export_to_markdown(report_month, comp_month, report_data, comparison_data, sc_data, comparison_sc_data, output_dir)
    
    print("\n" + "="*60)
    print("✅ テスト生成が完了しました！")
    print("="*60)

if __name__ == "__main__":
    main()

